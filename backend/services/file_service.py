"""
Serviço de geração de arquivos — PDF, DOCX, XLSX, PPTX.
Portado de netlify/functions/files.ts
"""

import io
import json
import logging
from typing import Tuple

logger = logging.getLogger(__name__)


def generate_pdf(title: str, content: str) -> bytes:
    """Gera PDF com reportlab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import cm

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                            topMargin=2*cm, bottomMargin=2*cm,
                            leftMargin=2*cm, rightMargin=2*cm)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "CustomTitle", parent=styles["Title"],
        fontSize=18, spaceAfter=20
    )
    body_style = ParagraphStyle(
        "CustomBody", parent=styles["Normal"],
        fontSize=12, leading=16
    )

    elements = [Paragraph(title, title_style), Spacer(1, 12)]

    for line in content.split("\n"):
        line = line.strip()
        if not line:
            elements.append(Spacer(1, 8))
        elif line.startswith("# "):
            elements.append(Paragraph(f"<b>{line[2:]}</b>", styles["Heading1"]))
        elif line.startswith("## "):
            elements.append(Paragraph(f"<b>{line[3:]}</b>", styles["Heading2"]))
        else:
            # Escape XML chars for reportlab
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            elements.append(Paragraph(safe, body_style))

    doc.build(elements)
    return buffer.getvalue()


def generate_docx(title: str, content: str) -> bytes:
    """Gera DOCX com python-docx."""
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)
    doc.add_paragraph("")  # spacer

    for line in content.split("\n"):
        trimmed = line.strip()
        if trimmed.startswith("# "):
            doc.add_heading(trimmed[2:], level=1)
        elif trimmed.startswith("## "):
            doc.add_heading(trimmed[3:], level=2)
        elif trimmed.startswith("### "):
            doc.add_heading(trimmed[4:], level=3)
        elif trimmed:
            p = doc.add_paragraph(trimmed)
            for run in p.runs:
                run.font.size = Pt(12)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def generate_xlsx(title: str, content: str) -> bytes:
    """Gera Excel com openpyxl."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] if title else "Dados"

    # Tentar parsear content como JSON (array de objetos ou array de arrays)
    try:
        data = json.loads(content)
        if isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], dict):
                # Array de objetos: headers das keys
                headers = list(data[0].keys())
                ws.append(headers)
                for row in data:
                    ws.append([str(row.get(h, "")) for h in headers])
            elif isinstance(data[0], list):
                # Array de arrays
                for row in data:
                    ws.append([str(cell) for cell in row])
            else:
                ws.append([title])
                ws.append([content])
        else:
            ws.append([title])
            ws.append([content])
    except (json.JSONDecodeError, TypeError):
        # Não é JSON, tratar como texto
        ws.append([title])
        for line in content.split("\n"):
            if line.strip():
                ws.append([line.strip()])

    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def generate_pptx(title: str, content: str) -> bytes:
    """Gera PPTX com python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Separar por marcador SLIDE:
    import re
    slides_content = re.split(r'SLIDE:', content, flags=re.IGNORECASE)
    slides_content = [s.strip() for s in slides_content if s.strip()]

    if not slides_content:
        # Fallback: um slide com tudo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    else:
        for slide_text in slides_content:
            lines = slide_text.split("\n")
            slide_title = lines[0].replace("*", "").replace("#", "").strip()
            slide_body = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_body

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# Mapa de tipos para geração
FILE_GENERATORS = {
    "pdf": {
        "func": generate_pdf,
        "mime": "application/pdf",
        "ext": "pdf",
    },
    "docx": {
        "func": generate_docx,
        "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "ext": "docx",
    },
    "excel": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "xlsx": {
        "func": generate_xlsx,
        "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "ext": "xlsx",
    },
    "pptx": {
        "func": generate_pptx,
        "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ext": "pptx",
    },
}


async def generate_file(file_type: str, title: str, content: str) -> Tuple[str, str]:
    """
    Gera arquivo e faz upload ao Supabase.

    Returns:
        (url_download, mensagem)
    """
    from backend.core.config import get_config
    from backend.core.supabase_client import upload_to_supabase

    config = get_config()
    file_type = file_type.lower()

    if file_type not in FILE_GENERATORS:
        raise ValueError(f"Tipo de arquivo inválido: {file_type}. Suportados: {list(FILE_GENERATORS.keys())}")

    gen = FILE_GENERATORS[file_type]

    # Gerar arquivo
    file_bytes = gen["func"](title, content)

    # Filename seguro
    import re
    safe_title = re.sub(r'[^a-z0-9]', '_', title.lower())[:50]
    filename = f"{safe_title}.{gen['ext']}"

    # Upload ao Supabase
    url = upload_to_supabase(
        bucket=config.supabase_storage_bucket,
        filename=filename,
        file_content=file_bytes,
        content_type=gen["mime"],
    )

    return url, f"{file_type.upper()} gerado com sucesso."
